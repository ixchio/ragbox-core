"""
Layer 2: AUTO-DOCUMENT INTELLIGENCE
MIME detection, smart parsing, and structural extraction.
"""
from abc import ABC, abstractmethod
from typing import Optional, Dict, Any, List
from pathlib import Path
import asyncio
from loguru import logger

from ragbox.models.documents import (
    Document, PDFDocument, CodeDocument, DocumentType
)

class BaseProcessor(ABC):
    @abstractmethod
    async def process(self, path: Path, file_hash: str) -> Optional[Document]:
        """Convert a raw file to a Document construct."""
        pass


class PDFProcessor(BaseProcessor):
    """PDF processor with robust text and basic table extraction."""
    async def process(self, path: Path, file_hash: str) -> Optional[PDFDocument]:
        try:
            import pdfplumber
            
            def _extract() -> PDFDocument:
                text_parts = []
                page_count = 0
                with pdfplumber.open(path) as pdf:
                    page_count = len(pdf.pages)
                    for i, page in enumerate(pdf.pages):
                        text_parts.append(f"--- Page {i+1} ---")
                        
                        # Extract text
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                            
                        # Extract tables (very basic stringification)
                        tables = page.extract_tables()
                        for table in tables:
                            text_parts.append("--- Table Data ---")
                            for row in table:
                                # Filter out Nones and join with visual separator
                                cleaned_row = [str(cell).strip() if cell else "" for cell in row]
                                text_parts.append(" | ".join(cleaned_row))
                                
                return PDFDocument(
                    id=file_hash,
                    path=path,
                    content="\n".join(text_parts),
                    page_count=page_count
                )
            return await asyncio.to_thread(_extract)
        except ImportError:
            logger.error("pdfplumber required for PDF processing. Install with: pip install pdfplumber")
            return None


class ImageProcessor(BaseProcessor):
    """Extracts text from images using OCR."""
    async def process(self, path: Path, file_hash: str) -> Optional[Document]:
        try:
            from paddleocr import PaddleOCR
            import logging
            
            def _extract() -> Document:
                # Suppress verbose paddleocr logging
                logging.getLogger("ppocr").setLevel(logging.ERROR)
                
                # Initialize OCR (use English model by default, limit to CPU to avoid CUDA setup overhead in general use)
                ocr = PaddleOCR(use_angle_cls=True, lang='en', use_gpu=False, show_log=False)
                result = ocr.ocr(str(path), cls=True)
                
                text_parts = []
                if result and result[0]: # result[0] is the list of lines for the first (and only) image
                     for line in result[0]:
                         # line is a tuple: (bounding_box, (text, confidence))
                         # e.g., ([[x,y], [x,y], [x,y], [x,y]], ('hello world', 0.99))
                         text, confidence = line[1]
                         if confidence > 0.6: # Configurable threshold later
                            text_parts.append(text)
                            
                return Document(
                    id=file_hash,
                    path=path,
                    content="\n".join(text_parts),
                    doc_type=DocumentType.TEXT
                )
            return await asyncio.to_thread(_extract)
        except ImportError:
            logger.error("paddleocr required for image processing. Install with: pip install paddleocr")
            return None


class PPTXProcessor(BaseProcessor):
    """PowerPoint (.pptx) processor — extracts text from all slides."""
    async def process(self, path: Path, file_hash: str) -> Optional[Document]:
        try:
            from pptx import Presentation

            def _extract() -> Document:
                prs = Presentation(str(path))
                slides_text = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_parts = [f"--- Slide {slide_num} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    slide_parts.append(text)
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_text = " | ".join(
                                    cell.text.strip() for cell in row.cells
                                )
                                if row_text.strip(" |"):
                                    slide_parts.append(row_text)
                    slides_text.append("\n".join(slide_parts))

                return Document(
                    id=file_hash,
                    path=path,
                    content="\n\n".join(slides_text),
                    doc_type=DocumentType.STRUCTURED,
                    metadata={"slide_count": len(prs.slides)},
                )
            return await asyncio.to_thread(_extract)
        except ImportError:
            logger.error("python-pptx required for PPTX processing. Install with: pip install python-pptx")
            return None


class CodeProcessor(BaseProcessor):
    async def process(self, path: Path, file_hash: str) -> Optional[CodeDocument]:
        def _read() -> CodeDocument:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return CodeDocument(
                id=file_hash,
                path=path,
                content=content,
                language=path.suffix.lstrip('.')
            )
        return await asyncio.to_thread(_read)


class TextProcessor(BaseProcessor):
    async def process(self, path: Path, file_hash: str) -> Optional[Document]:
        def _read() -> Document:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return Document(
                id=file_hash,
                path=path,
                content=content,
                doc_type=DocumentType.TEXT
            )
        return await asyncio.to_thread(_read)


class DocumentProcessorRouter:
    """Routes documents to correct processor automatically."""
    def __init__(self):
        self.processors = {
            # Documents
            ".pdf": PDFProcessor(),
            ".pptx": PPTXProcessor(),
            # Code
            ".py": CodeProcessor(),
            ".js": CodeProcessor(),
            ".ts": CodeProcessor(),
            ".tsx": CodeProcessor(),
            ".c": CodeProcessor(),
            ".cpp": CodeProcessor(),
            ".go": CodeProcessor(),
            ".rs": CodeProcessor(),
            ".java": CodeProcessor(),
            # Text / markup
            ".md": TextProcessor(),
            ".txt": TextProcessor(),
            ".csv": TextProcessor(),
            ".json": TextProcessor(),
            ".html": TextProcessor(),
            ".xml": TextProcessor(),
            # Images
            ".png": ImageProcessor(),
            ".jpg": ImageProcessor(),
            ".jpeg": ImageProcessor(),
        }
        self.fallback = TextProcessor()

    async def process(self, path: Path, file_hash: str) -> Optional[Document]:
        if not path.exists():
            return None
        ext = path.suffix.lower()
        processor = self.processors.get(ext, self.fallback)
        try:
            return await processor.process(path, file_hash)
        except Exception as e:
            logger.error(f"Failed to process {path} with {type(processor).__name__}: {e}")
            return None
